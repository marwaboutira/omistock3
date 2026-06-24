#!/usr/bin/env python3
"""
Script pour générer la présentation PowerPoint de soutenance OMISTOCK V3
17 slides avec effets dynamiques et design adapté aux logos
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap
import os

# Couleurs basées sur les logos OMISTOCK et ISTA
COLORS = {
    'emerald': RGBColor(16, 185, 129),
    'teal': RGBColor(20, 184, 166),
    'blue': RGBColor(59, 130, 246),
    'purple': RGBColor(139, 92, 246),
    'orange': RGBColor(249, 115, 22),
    'red': RGBColor(239, 68, 68),
    'green': RGBColor(34, 197, 94),
    'dark_bg': RGBColor(15, 23, 42),
    'light_bg': RGBColor(30, 41, 59),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(148, 163, 184),
    'accent': RGBColor(45, 212, 191)
}

def add_shape_with_animation(slide, shape_type, left, top, width, height, fill_color=None):
    """Ajoute une forme avec animation"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, font_color=COLORS['white'], bold=False, alignment=PP_ALIGN.LEFT):
    """Ajoute une zone de texte"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return text_box

def create_title_slide(prs):
    """Slide 1 : Page de Garde"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond sombre
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Logo OMISTOCK (simulé avec un cercle coloré)
    logo_left = Inches(3.5)
    logo_top = Inches(1.5)
    logo = add_shape_with_animation(slide, MSO_SHAPE.OVAL, logo_left, logo_top, Inches(3), Inches(1), COLORS['teal'])
    
    # Titre principal
    add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(1), 
                 "OMISTOCK V3", font_size=54, font_color=COLORS['teal'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # Sous-titre
    add_text_box(slide, Inches(1), Inches(4), Inches(8), Inches(0.5), 
                 "Système de Gestion de Stock Multi-Entreprise", font_size=24, font_color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    
    # Binôme
    add_text_box(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.5), 
                 "GACEM Sarra & BOUTIRA Marwa Assia", font_size=20, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    # Encadrant
    add_text_box(slide, Inches(1), Inches(6.2), Inches(8), Inches(0.5), 
                 "Encadrant : M. BENSAADA Ilies", font_size=18, font_color=COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    # Université
    add_text_box(slide, Inches(1), Inches(7), Inches(8), Inches(0.5), 
                 "ISTA / Université Oran 1 Ahmed Ben Bella", font_size=16, font_color=COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    # Notes
    slide.notes_slide.notes_text_frame.text = "Slide 1 : Page de Garde - Effet Zoom d'entrée"

def create_plan_slide(prs):
    """Slide 2 : Plan de la Présentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Plan de la Présentation", font_size=36, font_color=COLORS['teal'], bold=True)
    
    # Sommaire
    sections = [
        "Module 1 : Ouverture & Cadrage Professionnel",
        "Module 2 : Le Couvre-Feu Technologique (Problématique)",
        "Module 3 : Spécifications Métiers & Ingénierie de Conception",
        "Module 4 : Architecture Informatique & Diagrammes Systèmes",
        "Module 5 : Implémentation et Sécurité du Backend",
        "Module 6 : Démonstration en Conditions Réelles & Clôture"
    ]
    
    left = Inches(1)
    top = Inches(1.8)
    for i, section in enumerate(sections):
        color = COLORS['accent'] if i < 3 else COLORS['blue'] if i < 5 else COLORS['purple']
        add_text_box(slide, left, top, Inches(8), Inches(0.5), 
                     f"{i+1}. {section}", font_size=16, font_color=color, bold=True)
        top += Inches(0.6)
    
    slide.notes_slide.notes_text_frame.text = "Slide 2 : Plan - Effet Liste Révélée / Fragments"

def create_context_slide(prs):
    """Slide 3 : Contexte Logistique"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Contexte Logistique & Équilibre des Flux", font_size=36, font_color=COLORS['teal'], bold=True)
    
    # Deux boîtes pour les jauges
    box1 = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(3), COLORS['light_bg'])
    add_text_box(slide, Inches(0.7), Inches(2.2), Inches(3.6), Inches(0.5), 
                 "Niveau de Rupture", font_size=20, font_color=COLORS['red'], bold=True)
    add_text_box(slide, Inches(0.7), Inches(2.8), Inches(3.6), Inches(1.5), 
                 "• Perte de ventes\n• Clients frustrés\n• Image de marque dégradée", font_size=14, font_color=COLORS['white'])
    
    box2 = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2), Inches(4), Inches(3), COLORS['light_bg'])
    add_text_box(slide, Inches(5.2), Inches(2.2), Inches(3.6), Inches(0.5), 
                 "Niveau de Surstockage", font_size=20, font_color=COLORS['orange'], bold=True)
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(3.6), Inches(1.5), 
                 "• Dette de trésorerie\n• Coûts de stockage\n• Produits périmés", font_size=14, font_color=COLORS['white'])
    
    # Conclusion
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(1), 
                 "L'inventaire est un élément vital pour la survie économique des PME algériennes", 
                 font_size=18, font_color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 3 : Contexte - Effet Translation Latérale"

def create_limits_slide(prs):
    """Slide 4 : Limites des Systèmes Existants"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Limites des Systèmes Existants", font_size=36, font_color=COLORS['teal'], bold=True)
    
    # Trois panneaux d'alerte
    alerts = [
        ("Coût d'intégration SAP/Odoo", "Licences coûteuses\nFrais de conseil élevés\n> Budgets PME", COLORS['red']),
        ("Rigidité des systèmes", "Paramétrage complexe\nÉquipes IT dédiées\nMois de mise en place", COLORS['red']),
        ("Rejet des opérateurs", "Interface surchargée\nFonctionnalités inutiles\nProductivité réduite", COLORS['red'])
    ]
    
    left = Inches(0.5)
    for title, content, color in alerts:
        box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(3), Inches(3), COLORS['light_bg'])
        box.line.color.rgb = color
        box.line.width = Pt(3)
        add_text_box(slide, left + Inches(0.2), Inches(2.2), Inches(2.6), Inches(0.5), 
                     title, font_size=16, font_color=color, bold=True)
        add_text_box(slide, left + Inches(0.2), Inches(2.8), Inches(2.6), Inches(2), 
                     content, font_size=13, font_color=COLORS['white'])
        left += Inches(3.2)
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8), 
                 "Les PME abandonnent les ERP traditionnels après moins de 12 mois", 
                 font_size=18, font_color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 4 : Limites - Effet Pop-In Accumulatif"

def create_gap_slide(prs):
    """Slide 5 : Le Gap Agent-Ready"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Le Rupture Conceptuelle : Le \"Gap Agent-Ready\"", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Écran classique barré
    screen = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(2.5), COLORS['light_bg'])
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(2), Inches(0.5), 
                 "Interface Humaine", font_size=14, font_color=COLORS['gray'], alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(2), Inches(1), 
                 "Clics\nFormulaires\nDashboards", font_size=12, font_color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    
    # Croix rouge (X symbol)
    cross = add_shape_with_animation(slide, MSO_SHAPE.OVAL, Inches(2.3), Inches(2.8), Inches(0.6), Inches(0.6), COLORS['red'])
    add_text_box(slide, Inches(2.4), Inches(3), Inches(0.4), Inches(0.3), "✗", font_size=24, font_color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    
    # Flèche vers l'IA
    arrow = add_shape_with_animation(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(3), Inches(1.5), Inches(0.5), COLORS['accent'])
    
    # Processeur IA
    ai = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2), Inches(3), Inches(2.5), COLORS['blue'])
    add_text_box(slide, Inches(7), Inches(2.5), Inches(2), Inches(0.5), 
                 "Interface Agentique", font_size=14, font_color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7), Inches(3.2), Inches(2), Inches(1), 
                 "Sémantique\nAPI Standard\nAutonomie", font_size=12, font_color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(1), 
                 "Les logiciels actuels imposent une interface graphique humaine là où les machines ont besoin d'interfaces sémantiques", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 5 : Gap Agent-Ready - Effet Zoom Focus"

def create_mcp_slide(prs):
    """Slide 6 : Émergence des Écosystèmes Agentiques & MCP"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Émergence des Écosystèmes Agentiques & MCP", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # LLM
    llm = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(2.5), Inches(2), COLORS['purple'])
    add_text_box(slide, Inches(0.7), Inches(2.5), Inches(2.1), Inches(0.5), 
                 "LLM (Claude/GPT)", font_size=14, font_color=COLORS['white'], bold=True)
    add_text_box(slide, Inches(0.7), Inches(3.2), Inches(2.1), Inches(0.5), 
                 "Agent IA Autonome", font_size=12, font_color=COLORS['white'])
    
    # Connexions (simulées avec des rectangles fins)
    for i in range(3):
        conn = add_shape_with_animation(slide, MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(2.5 + i*0.5), Inches(2.5), Inches(0.05), COLORS['accent'])
    
    # MCP Server
    mcp = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(2), Inches(2.5), Inches(2), COLORS['teal'])
    add_text_box(slide, Inches(6.2), Inches(2.5), Inches(2.1), Inches(0.5), 
                 "Protocole MCP", font_size=14, font_color=COLORS['white'], bold=True)
    add_text_box(slide, Inches(6.2), Inches(3.2), Inches(2.1), Inches(0.5), 
                 "FastMCP Server", font_size=12, font_color=COLORS['white'])
    
    # Base de données
    db = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(4.5), Inches(2.5), Inches(1.5), COLORS['blue'])
    add_text_box(slide, Inches(6.2), Inches(5), Inches(2.1), Inches(0.5), 
                 "Base de Données", font_size=12, font_color=COLORS['white'])
    
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.8), 
                 "Protocole universel open-source MCP (Model Context Protocol) d'Anthropic - 2024", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 6 : MCP - Effet Révélation de Données"

def create_vision_slide(prs):
    """Slide 7 : La Vision OMISTOCK V3"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "La Vision OMISTOCK V3 : L'ERP Bi-Cible", font_size=36, font_color=COLORS['teal'], bold=True)
    
    # Colonne Humain
    human = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(3.5), COLORS['emerald'])
    add_text_box(slide, Inches(0.7), Inches(2.3), Inches(3.6), Inches(0.5), 
                 "👤 Interface Humaine", font_size=20, font_color=COLORS['white'], bold=True)
    add_text_box(slide, Inches(0.7), Inches(3), Inches(3.6), Inches(2), 
                 "• Web UI Responsive\n• PWA Mobile\n• Scanner QR Codes\n• Tableaux de bord", font_size=14, font_color=COLORS['white'])
    
    # Colonne Agent
    agent = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2), Inches(4), Inches(3.5), COLORS['purple'])
    add_text_box(slide, Inches(5.2), Inches(2.3), Inches(3.6), Inches(0.5), 
                 "🤖 Interface Agentique", font_size=20, font_color=COLORS['white'], bold=True)
    add_text_box(slide, Inches(5.2), Inches(3), Inches(3.6), Inches(2), 
                 "• FastMCP Server\n• API Sémantique\n• Outils MCP\n• Sécurité Scopée", font_size=14, font_color=COLORS['white'])
    
    add_text_box(slide, Inches(0.5), Inches(6), Inches(9), Inches(0.8), 
                 "Un système unique pour deux profils d'utilisateurs distincts", 
                 font_size=18, font_color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 7 : Vision - Effet Division 3D"

def create_algorithms_slide(prs):
    """Slide 8 : Algorithmes de Valorisation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Algorithmes de Valorisation et Logique Stock", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Trois cartes
    cards = [
        ("Coût Moyen Pondéré", "WAC = (Val. Init. + Coût Entrée) / (Qté Init. + Qté Entrée)", COLORS['emerald']),
        ("Point de Commande", "ROP = (Demande × Délai) + Stock Sécurité", COLORS['blue']),
        ("Règle FEFO", "ORDER BY date_expiration ASC", COLORS['purple'])
    ]
    
    left = Inches(0.5)
    for title, formula, color in cards:
        card = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(3), Inches(2.5), COLORS['light_bg'])
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_text_box(slide, left + Inches(0.2), Inches(2.2), Inches(2.6), Inches(0.5), 
                     title, font_size=14, font_color=color, bold=True)
        add_text_box(slide, left + Inches(0.2), Inches(2.8), Inches(2.6), Inches(1.5), 
                     formula, font_size=12, font_color=COLORS['white'])
        left += Inches(3.2)
    
    add_text_box(slide, Inches(0.5), Inches(5), Inches(9), Inches(1), 
                 "Les piliers mathématiques de l'application pour une gestion rigoureuse des stocks", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 8 : Algorithmes - Effet Rotation de Cartes"

def create_architecture_slide(prs):
    """Slide 9 : Architecture Générale 3-Tiers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Diagramme d'Architecture Générale 3-Tiers", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Frontend
    frontend = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(1), COLORS['emerald'])
    add_text_box(slide, Inches(1.5), Inches(2.3), Inches(7), Inches(0.5), 
                 "COUCHE CLIENT : Frontend (HTML5/Tailwind) + PWA Mobile", font_size=16, font_color=COLORS['white'], bold=True)
    
    # Flèche
    arrow1 = add_shape_with_animation(slide, MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(3.2), Inches(1), Inches(0.5), COLORS['accent'])
    
    # Backend
    backend = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(1), COLORS['blue'])
    add_text_box(slide, Inches(1.5), Inches(4.3), Inches(7), Inches(0.5), 
                 "SERVEUR APPLICATION : FastAPI (Python) + FastMCP Server", font_size=16, font_color=COLORS['white'], bold=True)
    
    # Flèche
    arrow2 = add_shape_with_animation(slide, MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(5.2), Inches(1), Inches(0.5), COLORS['accent'])
    
    # Database
    database = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6), Inches(8), Inches(1), COLORS['purple'])
    add_text_box(slide, Inches(1.5), Inches(6.3), Inches(7), Inches(0.5), 
                 "PERSISTANCE : SQLAlchemy ORM + SQLite Database", font_size=16, font_color=COLORS['white'], bold=True)
    
    slide.notes_slide.notes_text_frame.text = "Slide 9 : Architecture - Effet Plongée 3D"

def create_uml_slide(prs):
    """Slide 10 : Diagramme de Classes UML"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Modélisation des Données : Diagramme de Classes UML", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Tables
    tables = [
        ("Company", "id, name, created_at", COLORS['emerald']),
        ("User", "id, company_id, email, role", COLORS['blue']),
        ("Inventory", "id, company_id, product_id, quantity", COLORS['purple']),
        ("Movement", "id, inventory_id, type, quantity", COLORS['orange']),
        ("Batch", "id, inventory_id, expiry_date", COLORS['red'])
    ]
    
    left = Inches(0.5)
    top = Inches(2)
    for table_name, fields, color in tables:
        box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.5), Inches(1.2), COLORS['light_bg'])
        box.line.color.rgb = color
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(2.1), Inches(0.3), 
                     table_name, font_size=14, font_color=color, bold=True)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.5), Inches(2.1), Inches(0.6), 
                     fields, font_size=10, font_color=COLORS['white'])
        left += Inches(2.8)
        if left > Inches(6):
            left = Inches(0.5)
            top += Inches(1.5)
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8), 
                 "Multi-Tenancy : Étanchéité absolue via la clé étrangère globale company_id", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 10 : UML - Effet Zoom Haute Résolution"

def create_dfd_slide(prs):
    """Slide 11 : Diagramme de Flux de Données"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Diagramme de Flux de Données (DFD) — Entrée en Stock", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Étapes du flux
    steps = [
        ("Scan Marchandise", COLORS['emerald']),
        ("Validation Token JWT", COLORS['blue']),
        ("Vérification Capacité", COLORS['purple']),
        ("Recalcul WAC", COLORS['orange']),
        ("Écriture Base de Données", COLORS['red'])
    ]
    
    left = Inches(0.5)
    for step, color in steps:
        box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(1.8), Inches(1), COLORS['light_bg'])
        box.line.color.rgb = color
        add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(1.4), Inches(0.5), 
                     step, font_size=12, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
        if left < Inches(7):
            arrow = add_shape_with_animation(slide, MSO_SHAPE.RIGHT_ARROW, left + Inches(1.9), Inches(2.8), Inches(0.5), Inches(0.2), COLORS['accent'])
        left += Inches(2)
    
    add_text_box(slide, Inches(0.5), Inches(4), Inches(9), Inches(1.5), 
                 "Flux : Scan → Auth → Vérification → Calcul → Persistance\n\nValidation du token JWT, vérification de la capacité physique de l'entrepôt, et recalcul instantané de la valeur comptable du stock.", 
                 font_size=14, font_color=COLORS['white'])
    
    slide.notes_slide.notes_text_frame.text = "Slide 11 : DFD - Effet Flux Animé"

def create_sequence_slide(prs):
    """Slide 12 : Diagramme de Séquence Agentique"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Diagramme de Séquence de l'Interaction Agentique", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Lignes de vie
    actors = [
        ("Agent IA", COLORS['purple']),
        ("FastMCP", COLORS['teal']),
        ("Base de Données", COLORS['blue'])
    ]
    
    left = Inches(1)
    for actor, color in actors:
        # Ligne verticale (simulée avec un rectangle fin)
        line = add_shape_with_animation(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.8), Inches(2), Inches(0.05), Inches(3), COLORS['gray'])
        # Boîte acteur
        box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.6), Inches(0.5), color)
        add_text_box(slide, left + Inches(0.1), Inches(1.6), Inches(1.4), Inches(0.3), 
                     actor, font_size=12, font_color=COLORS['white'], alignment=PP_ALIGN.CENTER)
        left += Inches(3)
    
    # Messages
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(2), Inches(0.4), 
                 "get_stock_alerts()", font_size=10, font_color=COLORS['accent'])
    add_text_box(slide, Inches(4.5), Inches(3.5), Inches(2), Inches(0.4), 
                 "SELECT * FROM inventory", font_size=10, font_color=COLORS['accent'])
    add_text_box(slide, Inches(4.5), Inches(4.2), Inches(2), Inches(0.4), 
                 "Return data", font_size=10, font_color=COLORS['accent'])
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8), 
                 "L'IA interroge l'état critique du stock sans passer par l'interface visuelle", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 12 : Séquence - Effet Défilement Temporel"

def create_repository_slide(prs):
    """Slide 13 : Repository Pattern"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Implémentation du Patron de Conception \"Repository\"", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Fenêtre de code
    code_window = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(3.5), COLORS['light_bg'])
    code_window.line.color.rgb = COLORS['blue']
    code_window.line.width = Pt(2)
    
    # Code simulé
    code = """class InventoryRepository:
    def __init__(self, session):
        self.session = session
    
    def get_by_company(self, company_id):
        return self.session.query(Inventory)\\
            .filter(Inventory.company_id == company_id)\\
            .all()
    
    def update_quantity(self, inventory_id, qty):
        # Logique métier encapsulée
        inv = self.get_by_id(inventory_id)
        inv.quantity = qty
        self.session.commit()"""
    
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(3), 
                 code, font_size=11, font_color=COLORS['white'])
    
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(9), Inches(0.8), 
                 "Isolation de la logique d'accès aux données pour garantir une maintenance sans régression et des tests unitaires rapides", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 13 : Repository - Effet Fenêtre de Code"

def create_security_slide(prs):
    """Slide 14 : Politique de Sécurité"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Politique de Sécurité Hybride : JWT vs API-Keys", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Tableau comparatif
    # Humain
    human_box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(3), COLORS['emerald'])
    add_text_box(slide, Inches(0.7), Inches(2.2), Inches(3.6), Inches(0.5), 
                 "👤 Accès Humains", font_size=18, font_color=COLORS['white'], bold=True)
    add_text_box(slide, Inches(0.7), Inches(2.8), Inches(3.6), Inches(2), 
                 "• Authentification JWT Bearer\n• Cycle d'expiration court\n• Refresh Token sécurisé\n• Session HTTP Stateful", font_size=14, font_color=COLORS['white'])
    
    # Machine
    machine_box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2), Inches(4), Inches(3), COLORS['purple'])
    add_text_box(slide, Inches(5.2), Inches(2.2), Inches(3.6), Inches(0.5), 
                 "🤖 Accès Machines", font_size=18, font_color=COLORS['white'], bold=True)
    add_text_box(slide, Inches(5.2), Inches(2.8), Inches(3.6), Inches(2), 
                 "• Clés X-API-Key révocables\n• Scopes précis (read/write)\n• Pas de session stateless\n• Audit trail complet", font_size=14, font_color=COLORS['white'])
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8), 
                 "Sécurité adaptée au profil d'utilisateur pour une protection maximale", 
                 font_size=16, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 14 : Sécurité - Effet Verrou Cyber"

def create_demo_slide(prs):
    """Slide 15 : Scénario de Simulation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Scénario de Simulation : Le Réapprovisionnement Autonome", font_size=32, font_color=COLORS['teal'], bold=True)
    
    # Console simulée
    console = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(3), COLORS['light_bg'])
    console.line.color.rgb = COLORS['green']
    console.line.width = Pt(2)
    
    logs = """[2024-06-15 10:23:45] INFO: Agent IA initialized
[2024-06-15 10:23:46] INFO: Checking stock alerts...
[2024-06-15 10:23:47] WARNING: Product SKU-001 below ROP (qty: 5, rop: 20)
[2024-06-15 10:23:48] INFO: Calculating EOQ for SKU-001...
[2024-06-15 10:23:49] INFO: EOQ = 150 units (Wilson formula)
[2024-06-15 10:23:50] INFO: Creating restock proposal...
[2024-06-15 10:23:51] INFO: Proposal ID #1234 created - PENDING approval
[2024-06-15 10:23:52] SUCCESS: Restock proposal sent to human operator"""
    
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(2.5), 
                 logs, font_size=10, font_color=COLORS['emerald'])
    
    add_text_box(slide, Inches(0.5), Inches(5.3), Inches(9), Inches(1), 
                 "Démonstration de la proactivité de la solution OMISTOCK V3 : L'agent IA détecte et propose automatiquement les actions de réapprovisionnement", 
                 font_size=14, font_color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 15 : Demo - Effet Console Live"

def create_bilan_slide(prs):
    """Slide 16 : Bilan Technique"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Bilan Technique & Objectifs Atteints", font_size=36, font_color=COLORS['teal'], bold=True)
    
    # Check-list
    objectives = [
        "✓ Système multi-tenant opérationnel",
        "✓ API REST performante (FastAPI)",
        "✓ Passerelle sémantique MCP validée",
        "✓ Algorithmes WAC/ROP/FEFO implémentés",
        "✓ Interface PWA responsive",
        "✓ Sécurité hybride JWT/API-Keys",
        "✓ Déploiement Cloud (Render)",
        "✓ Auditabilité complète (SHA-256)"
    ]
    
    left = Inches(2)
    top = Inches(2)
    for obj in objectives:
        add_text_box(slide, left, top, Inches(6), Inches(0.4), 
                     obj, font_size=16, font_color=COLORS['emerald'], bold=True)
        top += Inches(0.5)
    
    add_text_box(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.8), 
                 "Synthèse de la valeur ajoutée apportée par le projet de fin d'études", 
                 font_size=18, font_color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 16 : Bilan - Effet Check-list Animée"

def create_perspectives_slide(prs):
    """Slide 17 : Perspectives & Clôture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark_bg']
    
    # Titre
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 
                 "Perspectives d'Avenir & Clôture", font_size=36, font_color=COLORS['teal'], bold=True)
    
    # Perspectives
    perspectives = [
        ("Migration PostgreSQL", "Pour gérer un niveau de concurrence élevé"),
        ("Jetons JIT éphémères", "Sécurisation renforcée des agents IA"),
        ("Webhooks en mode Push", "Alertes temps réel vers les agents"),
        ("Algorithmes prédictifs", "Deep-learning pour la demande"),
        ("Déploiement Kubernetes", "Scalabilité production")
    ]
    
    left = Inches(0.5)
    top = Inches(2)
    for title, desc in perspectives:
        box = add_shape_with_animation(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4), Inches(0.8), COLORS['light_bg'])
        box.line.color.rgb = COLORS['blue']
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(3.6), Inches(0.3), 
                     title, font_size=14, font_color=COLORS['blue'], bold=True)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.4), Inches(3.6), Inches(0.3), 
                     desc, font_size=11, font_color=COLORS['white'])
        left += Inches(4.5)
        if left > Inches(5):
            left = Inches(0.5)
            top += Inches(1)
    
    # Remerciements
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(1), 
                 "Nous vous remercions pour votre attention et sommes prêtes pour vos questions", 
                 font_size=20, font_color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Slide 17 : Perspectives - Effet Zoom Arrière Global"

def main():
    """Fonction principale"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Créer les 17 slides
    create_title_slide(prs)
    create_plan_slide(prs)
    create_context_slide(prs)
    create_limits_slide(prs)
    create_gap_slide(prs)
    create_mcp_slide(prs)
    create_vision_slide(prs)
    create_algorithms_slide(prs)
    create_architecture_slide(prs)
    create_uml_slide(prs)
    create_dfd_slide(prs)
    create_sequence_slide(prs)
    create_repository_slide(prs)
    create_security_slide(prs)
    create_demo_slide(prs)
    create_bilan_slide(prs)
    create_perspectives_slide(prs)
    
    # Sauvegarder
    output_file = "soutenance_omistock_v3.pptx"
    prs.save(output_file)
    print(f"Présentation sauvegardée : {output_file}")
    print(f"Nombre de slides : {len(prs.slides)}")

if __name__ == "__main__":
    main()
