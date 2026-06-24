#!/usr/bin/env python3
"""
Script pour convertir la présentation HTML en PowerPoint (PPTX)
Utilise la bibliothèque python-pptx pour générer un fichier modifiable
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re

def extract_slide_content(html_file):
    """Extrait le contenu de chaque slide depuis le fichier HTML"""
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = []
    
    # Trouver toutes les slides
    slide_divs = soup.find_all('div', class_='slide')
    
    for slide_div in slide_divs:
        slide_data = {
            'title': '',
            'subtitle': '',
            'content': [],
            'script': ''
        }
        
        # Extraire le titre principal
        title_elem = slide_div.find('h2')
        if title_elem:
            slide_data['title'] = title_elem.get_text(strip=True) or title_elem.get_text()
        
        # Extraire le sous-titre
        subtitle_elem = slide_div.find('p')
        if subtitle_elem and slide_div.find('h2'):
            slide_data['subtitle'] = subtitle_elem.get_text(strip=True)
        
        # Extraire le contenu des cartes
        cards = slide_div.find_all('div', class_='glass-card')
        for card in cards:
            card_content = []
            
            # Titre de la carte
            card_title = card.find(['h3', 'h4'])
            if card_title:
                card_content.append(('title', card_title.get_text(strip=True)))
            
            # Texte de la carte
            card_text = card.find('p')
            if card_text:
                card_content.append(('text', card_text.get_text(strip=True)))
            
            # Code/formule
            code_elem = card.find('div', class_=re.compile(r'bg-|font-mono'))
            if code_elem:
                card_content.append(('code', code_elem.get_text(strip=True)))
            
            if card_content:
                slide_data['content'].append(card_content)
        
        # Extraire le script oral
        script_elem = slide_div.find('div', class_='script-box')
        if script_elem:
            script_text = script_elem.get_text(strip=True)
            # Enlever "Script Oral :" du début
            if 'Script Oral :' in script_text:
                script_text = script_text.split('Script Oral :')[1].strip()
            slide_data['script'] = script_text
        
        slides.append(slide_data)
    
    return slides

def create_presentation(slides_data, output_file):
    """Crée la présentation PowerPoint"""
    prs = Presentation()
    
    # Définir les couleurs (basées sur le design HTML)
    colors = {
        'emerald': RGBColor(16, 185, 129),
        'blue': RGBColor(59, 130, 246),
        'purple': RGBColor(139, 92, 246),
        'orange': RGBColor(249, 115, 22),
        'dark_bg': RGBColor(15, 23, 42),
        'light_text': RGBColor(255, 255, 255),
        'accent': RGBColor(20, 184, 166)
    }
    
    for i, slide_data in enumerate(slides_data):
        # Choisir le layout selon le contenu
        if i == 0:
            # Première slide - layout titre
            slide_layout = prs.slide_layouts[0]  # Title slide
        elif len(slide_data['content']) <= 2:
            slide_layout = prs.slide_layouts[1]  # Title and Content
        else:
            slide_layout = prs.slide_layouts[5]  # Blank pour layout personnalisé
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Fond de slide sombre
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['dark_bg']
        
        # Ajouter le titre
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.color.rgb = colors['light_text']
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Ajouter le sous-titre si présent
        if slide_data['subtitle']:
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_data['subtitle']
                subtitle.text_frame.paragraphs[0].font.size = Pt(18)
                subtitle.text_frame.paragraphs[0].font.color.rgb = colors['accent']
                subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Ajouter le contenu des cartes
        if slide_data['content']:
            # Positionnement pour le contenu
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(9)
            height = Inches(4)
            
            # Créer une zone de texte pour le contenu
            content_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for card_idx, card_content in enumerate(slide_data['content']):
                if card_idx > 0:
                    p = text_frame.add_paragraph()
                    p.space_before = Pt(12)
                
                for elem_type, elem_text in card_content:
                    if elem_type == 'title':
                        p = text_frame.add_paragraph()
                        p.text = f"• {elem_text}"
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.font.color.rgb = colors['accent']
                        p.space_after = Pt(6)
                    elif elem_type == 'text':
                        p = text_frame.add_paragraph()
                        p.text = f"  {elem_text}"
                        p.font.size = Pt(14)
                        p.font.color.rgb = colors['light_text']
                        p.level = 1
                        p.space_after = Pt(6)
                    elif elem_type == 'code':
                        p = text_frame.add_paragraph()
                        p.text = f"  Code: {elem_text}"
                        p.font.size = Pt(12)
                        p.font.color.rgb = colors['blue']
                        p.font.italic = True
                        p.level = 1
                        p.space_after = Pt(6)
        
        # Ajouter le script oral en bas (optionnel - en notes)
        if slide_data['script']:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = f"Script Oral: {slide_data['script']}"
    
    # Sauvegarder la présentation
    prs.save(output_file)
    print(f"Présentation sauvegardée: {output_file}")

if __name__ == "__main__":
    html_file = "soutenance.html"
    output_file = "soutenance.pptx"
    
    print("Extraction du contenu HTML...")
    slides_data = extract_slide_content(html_file)
    print(f"{len(slides_data)} slides extraites")
    
    print("Génération du PowerPoint...")
    create_presentation(slides_data, output_file)
    
    print("Conversion terminée!")
